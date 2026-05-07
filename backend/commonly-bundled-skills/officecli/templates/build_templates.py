"""
Build Commonly-branded officecli starter templates with {{key}} placeholders.

Outputs 3 files into backend/commonly-bundled-skills/officecli/templates/.
Agents merge them with a JSON data object via:

  officecli merge templates/commonly-brief.docx out.docx --data '{"title":...}'

This produces a fully-styled, fully-populated deliverable in one shot —
instead of `create` + 30-50 individual `officecli set` calls.

Palette (from frontend/src/v2/v2.css):
  --v2-accent          #2f6feb
  --v2-accent-strong   #1f55c9
  --v2-accent-soft     #e8efff
  --v2-accent-deep     #14306f
  --v2-text-primary    #111827
  --v2-text-secondary  #4b5563
"""

import os
from docx import Document
from docx.shared import Pt, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH

from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter

from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt, Emu
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.shapes import MSO_SHAPE

OUT = "backend/commonly-bundled-skills/officecli/templates"
os.makedirs(OUT, exist_ok=True)

ACCENT = "2f6feb"
ACCENT_STRONG = "1f55c9"
ACCENT_SOFT = "e8efff"
ACCENT_DEEP = "14306f"
TEXT_PRIMARY = "111827"
TEXT_SECONDARY = "4b5563"


def build_docx():
    """commonly-brief.docx — title + subtitle + 3 sections.

    Placeholders consumed by `officecli merge --data`:
      title, subtitle, context, finding_1, finding_2, finding_3,
      recommendation, footer
    """
    doc = Document()

    section = doc.sections[0]
    section.top_margin = Inches(1)
    section.bottom_margin = Inches(1)
    section.left_margin = Inches(1)
    section.right_margin = Inches(1)

    # Body default
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(11)
    style.font.color.rgb = RGBColor.from_string(TEXT_PRIMARY)

    # H1: deep accent, 24pt
    h1 = doc.styles["Heading 1"]
    h1.font.name = "Calibri"
    h1.font.size = Pt(24)
    h1.font.bold = True
    h1.font.color.rgb = RGBColor.from_string(ACCENT_DEEP)

    # H2: primary text, 16pt
    h2 = doc.styles["Heading 2"]
    h2.font.name = "Calibri"
    h2.font.size = Pt(16)
    h2.font.bold = True
    h2.font.color.rgb = RGBColor.from_string(TEXT_PRIMARY)

    # Title (placeholder)
    t = doc.add_paragraph(style="Heading 1")
    t.add_run("{{title}}")

    # Subtitle (placeholder; italic, secondary)
    sub = doc.add_paragraph()
    sub_run = sub.add_run("{{subtitle}}")
    sub_run.font.size = Pt(12)
    sub_run.font.color.rgb = RGBColor.from_string(TEXT_SECONDARY)
    sub_run.italic = True

    doc.add_paragraph()

    # Context section
    doc.add_paragraph("Context", style="Heading 2")
    doc.add_paragraph("{{context}}")

    # Findings (3 bulleted placeholders)
    doc.add_paragraph("Findings", style="Heading 2")
    for key in ("finding_1", "finding_2", "finding_3"):
        doc.add_paragraph(f"• {{{{{key}}}}}")

    # Recommendation
    doc.add_paragraph("Recommendation", style="Heading 2")
    doc.add_paragraph("{{recommendation}}")

    # Footer
    fp = section.footer.paragraphs[0]
    fp.text = "{{footer}}"
    fp.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    fp.runs[0].font.size = Pt(9)
    fp.runs[0].font.color.rgb = RGBColor.from_string(TEXT_SECONDARY)

    out = os.path.join(OUT, "commonly-brief.docx")
    doc.save(out)
    print(f"✓ docx → {out} ({os.path.getsize(out)} bytes)")
    return out


def build_xlsx():
    """commonly-data.xlsx — 5-column data matrix with styled header.

    The merge command substitutes {{key}} in cell values. For arbitrary
    row counts, the agent should populate the template directly via
    `officecli add` for rows beyond row 2 — the styled header + frozen
    pane + column widths are the value here.

    Placeholders in row 2 (one sample row):
      sample_item, sample_status, sample_owner, sample_notes, sample_updated
    """
    wb = Workbook()
    ws = wb.active
    ws.title = "Data"

    headers = ["Item", "Status", "Owner", "Notes", "Updated"]
    for col_idx, h in enumerate(headers, 1):
        cell = ws.cell(row=1, column=col_idx, value=h)
        cell.font = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
        cell.fill = PatternFill("solid", fgColor=ACCENT_STRONG)
        cell.alignment = Alignment(horizontal="left", vertical="center", indent=1)

    # One sample row with placeholders
    sample_keys = ["sample_item", "sample_status", "sample_owner",
                   "sample_notes", "sample_updated"]
    for col_idx, key in enumerate(sample_keys, 1):
        cell = ws.cell(row=2, column=col_idx, value="{{" + key + "}}")
        cell.font = Font(name="Calibri", size=10, color=TEXT_PRIMARY)
        cell.alignment = Alignment(vertical="center", indent=1)
        cell.border = Border(bottom=Side(style="thin", color="e5e7eb"))

    widths = [22, 14, 14, 40, 14]
    for col_idx, w in enumerate(widths, 1):
        ws.column_dimensions[get_column_letter(col_idx)].width = w

    ws.freeze_panes = "A2"

    out = os.path.join(OUT, "commonly-data.xlsx")
    wb.save(out)
    print(f"✓ xlsx → {out} ({os.path.getsize(out)} bytes)")
    return out


def build_pptx():
    """commonly-deck.pptx — 3-slide branded deck (cover / content / closing).

    Placeholders:
      title, subtitle  (cover)
      slide2_heading, slide2_p1, slide2_p2, slide2_p3  (content)
      closing_heading, closing_body, closing_meta  (closing)
    """
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)
    blank = prs.slide_layouts[6]

    accent_rgb = PptRGB(0x2f, 0x6f, 0xeb)
    deep_rgb = PptRGB(0x14, 0x30, 0x6f)
    primary_rgb = PptRGB(0x11, 0x18, 0x27)
    secondary_rgb = PptRGB(0x4b, 0x55, 0x63)
    muted_rgb = PptRGB(0x7b, 0x84, 0x94)

    # ---- Cover slide ----
    s1 = prs.slides.add_slide(blank)

    bar = s1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0),
        prs.slide_width, PptInches(0.5)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_rgb
    bar.line.fill.background()

    title_box = s1.shapes.add_textbox(
        PptInches(0.7), PptInches(2.6), PptInches(12), PptInches(1.5)
    )
    title_p = title_box.text_frame.paragraphs[0]
    title_p.text = "{{title}}"
    title_p.font.name = "Calibri"
    title_p.font.size = PptPt(48)
    title_p.font.bold = True
    title_p.font.color.rgb = deep_rgb

    sub_box = s1.shapes.add_textbox(
        PptInches(0.7), PptInches(4.2), PptInches(12), PptInches(0.8)
    )
    sub_p = sub_box.text_frame.paragraphs[0]
    sub_p.text = "{{subtitle}}"
    sub_p.font.name = "Calibri"
    sub_p.font.size = PptPt(20)
    sub_p.font.color.rgb = secondary_rgb

    foot_box = s1.shapes.add_textbox(
        PptInches(0.7), PptInches(7.0), PptInches(12), PptInches(0.4)
    )
    foot_p = foot_box.text_frame.paragraphs[0]
    foot_p.text = "Commonly"
    foot_p.font.name = "Calibri"
    foot_p.font.size = PptPt(11)
    foot_p.font.color.rgb = muted_rgb

    # ---- Content slide ----
    s2 = prs.slides.add_slide(blank)

    bar2 = s2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0),
        prs.slide_width, PptInches(0.25)
    )
    bar2.fill.solid()
    bar2.fill.fore_color.rgb = accent_rgb
    bar2.line.fill.background()

    h_box = s2.shapes.add_textbox(
        PptInches(0.7), PptInches(0.6), PptInches(12), PptInches(0.9)
    )
    h_p = h_box.text_frame.paragraphs[0]
    h_p.text = "{{slide2_heading}}"
    h_p.font.name = "Calibri"
    h_p.font.size = PptPt(32)
    h_p.font.bold = True
    h_p.font.color.rgb = primary_rgb

    body_box = s2.shapes.add_textbox(
        PptInches(0.7), PptInches(1.8), PptInches(12), PptInches(5)
    )
    tf = body_box.text_frame
    tf.word_wrap = True
    for i, key in enumerate(("slide2_p1", "slide2_p2", "slide2_p3")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + "{{" + key + "}}"
        p.font.name = "Calibri"
        p.font.size = PptPt(20)
        p.font.color.rgb = primary_rgb
        p.space_after = PptPt(12)

    pnum = s2.shapes.add_textbox(
        PptInches(12.2), PptInches(7.0), PptInches(0.7), PptInches(0.3)
    )
    pnum_p = pnum.text_frame.paragraphs[0]
    pnum_p.text = "2"
    pnum_p.font.name = "Calibri"
    pnum_p.font.size = PptPt(11)
    pnum_p.font.color.rgb = muted_rgb

    # ---- Closing slide ----
    s3 = prs.slides.add_slide(blank)

    panel = s3.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0),
        PptInches(4.5), prs.slide_height
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = deep_rgb
    panel.line.fill.background()

    pt_box = s3.shapes.add_textbox(
        PptInches(0.5), PptInches(3.2), PptInches(4), PptInches(1)
    )
    pt_p = pt_box.text_frame.paragraphs[0]
    pt_p.text = "{{closing_heading}}"
    pt_p.font.name = "Calibri"
    pt_p.font.size = PptPt(36)
    pt_p.font.bold = True
    pt_p.font.color.rgb = PptRGB(0xff, 0xff, 0xff)

    rc_box = s3.shapes.add_textbox(
        PptInches(5.2), PptInches(2.6), PptInches(7.5), PptInches(3)
    )
    rc_tf = rc_box.text_frame
    rc_tf.word_wrap = True
    rc_p1 = rc_tf.paragraphs[0]
    rc_p1.text = "{{closing_body}}"
    rc_p1.font.name = "Calibri"
    rc_p1.font.size = PptPt(28)
    rc_p1.font.color.rgb = primary_rgb
    rc_tf.add_paragraph().text = ""
    rc_p3 = rc_tf.add_paragraph()
    rc_p3.text = "{{closing_meta}}"
    rc_p3.font.name = "Calibri"
    rc_p3.font.size = PptPt(16)
    rc_p3.font.color.rgb = secondary_rgb

    out = os.path.join(OUT, "commonly-deck.pptx")
    prs.save(out)
    print(f"✓ pptx → {out} ({os.path.getsize(out)} bytes)")
    return out


if __name__ == "__main__":
    build_docx()
    build_xlsx()
    build_pptx()
